import os.path
import glob
from pptx import Presentation
from pptx.util import Cm  # 导入单位-厘米
from pptx.util import Pt  # 导入单位
from pptx.dml.color import RGBColor


def insert_images(pngfile):
    img_path = pngfile
    img_left, img_top, img_width, img_height = Cm(11.7), Cm(8.53), Cm(2), Cm(2)  # 设置图片位置和大小，距离左边、上边，图片的宽、高
    pic = slide.shapes.add_picture(img_path, img_left, img_top, img_width, img_height)
    tx_left, tx_top, tx_width, tx_height = Cm(17), Cm(8.53), Cm(2), Cm(1)  # 设置文本框的位置和大小，距离左边、上边，文本框的宽、高
    textbox = slide.shapes.add_textbox(tx_left, tx_top, tx_width, tx_height)
    text_frame = textbox.text_frame
    p1 = text_frame.paragraphs[0]  # 文本框建好之后已经默认有了第一段，即paragraphs[0]
    p1.text = os.path.basename(os.path.splitext(pngfile)[0])  # paragraphs[0]的文字内容是png格式图片的文件名
    font = p1.font  # 为文字指定字体格式
    font.name = 'Calibre'  # 设置字体
    font.size = Pt(20)  # 设置文字大小（此处为磅）
    font.color.rgb = RGBColor(0, 0, 0)  # 设置文字颜色（注意开头需要import RGBColor
    # font.bold = True 文字加粗
    # font.italic = True 文字斜体
    # font.underline = True 文字下划线
    p2 = text_frame.add_paragraph()  # 在文本框中添加新段落
    p2.text = '祝您春节快乐'


prs = Presentation()
for pngfile in glob.glob("C:/Users/Willem/Jupyter_Exercise/已处理的图片/*.png"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 插入图片之前需要先插入一页空白PPT
    insert_images(pngfile)
prs.save('test.pptx')  # 保存PPT文件
print('完成')
