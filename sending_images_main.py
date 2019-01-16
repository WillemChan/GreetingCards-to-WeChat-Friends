import itchat
import os
import os.path
import math
import glob
from pptx import Presentation
from pptx.util import Cm  # 导入单位-厘米
from pptx.util import Pt  # 导入单位
from pptx.dml.color import RGBColor
from PIL import Image
from Download_HeadImages import dld_hImages
from Reshape_Images_to_Circle import reshape_image_to_circle


def get_PPTx(images_bef, images_aft, pptx_bef, pptx_aft):  # 分别代表头像源文件保存路径、剪切为圆形的头像保存路径、需要批量添加图片和文字的PPT文件路径以及导出的PPT文件路径
    dld_hImages(images_bef)  # 调用Download_HeadImages中的函数，用来批量下载好友头像
    print('头像下载完毕')
    os.mkdir(images_aft)  # 新建文件夹用于存放剪切为圆形的头像文件
    os.chdir(images_aft)  # 将工作目录指定为images_aft文件夹
    for jpgfile in glob.glob(images_bef + '/*.jpg'):  # 保存的头像文件都是.jpg格式
        reshape_image_to_circle(jpgfile, images_aft)  # 将处理后的头像文件保存到images_aft路径下
    print('图片处理完毕')

    def insert_images(pngfile):  # 由于涉及到slide变量，所以需要嵌套函数，该函数也是Insert_Images_into_pptx中的函数
        img_path = pngfile
        img_left, img_top, img_width, img_height = Cm(11.7), Cm(8.53), Cm(2), Cm(2)  # 设置图片位置和大小，距离左边、上边，图片的宽、高
        pic = slide.shapes.add_picture(img_path, img_left, img_top, img_width, img_height)
        tx_left, tx_top, tx_width, tx_height = Cm(17), Cm(8.53), Cm(2), Cm(1)  # 设置文本框的位置和大小，距离左边、上边，文本框的宽、高
        textbox = slide.shapes.add_textbox(tx_left, tx_top, tx_width, tx_height)
        text_frame = textbox.text_frame
        p1 = text_frame.paragraphs[0]  # 文本框建好之后已经默认有了第一段，即paragraphs[0]
        p1.text = os.path.basename(os.path.splitext(pngfile)[0])  # paragraphs[0]的文字内容是png格式图片的文件名
        font = p1.font  # 为文字指定字体格式
        font.name = 'Calibre'  # 设置字体
        font.size = Pt(20)  # 设置文字大小（此处为磅）
        font.color.rgb = RGBColor(0, 0, 0)  # 设置文字颜色（注意开头需要import RGBColor
        # font.bold = True 文字加粗
        # font.italic = True 文字斜体
        # font.underline = True 文字下划线
        p2 = text_frame.add_paragraph()  # 在文本框中添加新段落
        p2.text = '祝您春节快乐'

    prs = Presentation(pptx_bef)
    pages = 0  # 需要对未处理的PPT.pptx中的所有页面插入图片
    for slide, pngfile in zip(prs.slides, glob.glob(images_aft + '/*.png')):
        if pages > len(prs.slides) - 1:
            continue
        insert_images(pngfile)
        pages += 1
    prs.save(pptx_aft)  # 保存PPT文件
    print('PPT处理完毕')


if __name__ == "__main__":
    images_bef = 'C:/Users/headImages_bef'  # 头像源文件保存路径
    images_aft = 'C:/Users/headImages_aft'  # 剪切为圆形的头像保存路径
    pptx_bef = 'C:/Users/未处理的PPT.pptx'  # 批量添加图片和文字的PPT文件路径
    pptx_aft = 'C:/Users/已处理的PPT.pptx'  # 导出的PPT文件路径
    get_PPTx(images_bef, images_aft, pptx_bef, pptx_aft)
